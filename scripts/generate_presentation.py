# -*- coding: utf-8 -*-
"""Generate final presentation PPT for 银发健康助手."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent.parent / "银发健康助手-期末汇报PPT-v3.pptx"

FONT_NAME = "微软雅黑"
MIN_SIZE = 20

# Warm & light palette — elderly-friendly, not tech/AI look
COVER_BG = RGBColor(0xFF, 0xF5, 0xEB)          # warm cream
COVER_DECOR_1 = RGBColor(0xF5, 0xD0, 0xC5)      # soft peach
COVER_DECOR_2 = RGBColor(0xD4, 0xE8, 0xC8)      # soft sage
HEADER = RGBColor(0xE8, 0xC4, 0xA8)             # warm apricot sand
HEADER_LINE = RGBColor(0xB8, 0xD4, 0xA8)        # soft sage line
ACCENT = RGBColor(0x6B, 0x9B, 0x6E)             # gentle sage green
ACCENT_WARM = RGBColor(0xC9, 0x76, 0x5D)        # soft terracotta
TITLE = RGBColor(0x5C, 0x40, 0x33)              # warm brown
TEXT = RGBColor(0x4A, 0x40, 0x35)               # warm charcoal
SUBTEXT = RGBColor(0x7A, 0x6E, 0x63)            # warm gray-brown
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PAGE_BG = RGBColor(0xFF, 0xFB, 0xF7)            # warm white
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xED, 0xD9, 0xC8)        # warm beige border
HIGHLIGHT_BG = RGBColor(0xF5, 0xF9, 0xF0)       # soft mint cream
TABLE_ALT = RGBColor(0xFD, 0xF6, 0xEF)          # warm cream row
TABLE_HEADER = RGBColor(0xF0, 0xDB, 0xC8)       # light peach header


def style_font(font, size=MIN_SIZE, bold=False, color=TEXT, italic=False):
    font.name = FONT_NAME
    font.size = Pt(max(size, MIN_SIZE))
    font.bold = bold
    font.italic = italic
    if color is not None:
        font.color.rgb = color
    rPr = font._element
    rPr.set(qn("a:latin"), FONT_NAME)
    rPr.set(qn("a:ea"), FONT_NAME)
    rPr.set(qn("a:cs"), FONT_NAME)


def style_paragraph(p, size=MIN_SIZE, bold=False, color=TEXT, italic=False):
    if p.runs:
        for run in p.runs:
            style_font(run.font, size, bold, color, italic)
    else:
        style_font(p.font, size, bold, color, italic)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, title_text):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = HEADER
    bar.line.fill.background()

    line = slide.shapes.add_shape(1, Inches(0), Inches(0.95), Inches(10), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = HEADER_LINE
    line.line.fill.background()

    tf = bar.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    style_paragraph(p, size=28, bold=True, color=TITLE)
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.55)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_bullets(slide, items, left=0.55, top=1.25, width=8.9, height=5.6, font_size=22):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, sub) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        style_paragraph(p, size=font_size, bold=("：" in text and "•" not in text), color=TEXT)
        p.space_after = Pt(8)
        if sub:
            sp = tf.add_paragraph()
            sp.text = sub
            style_paragraph(sp, size=MIN_SIZE, color=SUBTEXT)
            sp.level = 1
            sp.space_after = Pt(12)


def add_demo_card(slide, y, title, user_says, system_does):
    card = slide.shapes.add_shape(1, Inches(0.45), Inches(y), Inches(9.1), Inches(1.72))
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(1.5)

    left_bar = slide.shapes.add_shape(1, Inches(0.45), Inches(y), Inches(0.08), Inches(1.72))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = ACCENT
    left_bar.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.65), Inches(y + 0.1), Inches(8.7), Inches(1.55))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = title
    style_paragraph(p0, size=22, bold=True, color=ACCENT)

    p1 = tf.add_paragraph()
    p1.text = f'用户说："{user_says}"'
    style_paragraph(p1, size=MIN_SIZE, color=TEXT)
    p1.space_before = Pt(6)

    p2 = tf.add_paragraph()
    p2.text = f"系统做：{system_does}"
    style_paragraph(p2, size=MIN_SIZE, color=SUBTEXT)
    p2.space_before = Pt(4)


def add_warm_decor(slide):
    """Soft organic shapes — warm community feel, not tech circles."""
    blob1 = slide.shapes.add_shape(9, Inches(7.0), Inches(-0.5), Inches(3.2), Inches(3.2))
    blob1.fill.solid()
    blob1.fill.fore_color.rgb = COVER_DECOR_1
    blob1.fill.transparency = 0.55
    blob1.line.fill.background()

    blob2 = slide.shapes.add_shape(9, Inches(-0.4), Inches(5.8), Inches(2.6), Inches(2.6))
    blob2.fill.solid()
    blob2.fill.fore_color.rgb = COVER_DECOR_2
    blob2.fill.transparency = 0.5
    blob2.line.fill.background()

    blob3 = slide.shapes.add_shape(9, Inches(8.2), Inches(5.5), Inches(1.8), Inches(1.8))
    blob3.fill.solid()
    blob3.fill.fore_color.rgb = COVER_DECOR_2
    blob3.fill.transparency = 0.65
    blob3.line.fill.background()


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- Slide 1: Cover (light warm) ---
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1, COVER_BG)
    add_warm_decor(s1)

    title_box = s1.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(8.4), Inches(1.3))
    tp = title_box.text_frame.paragraphs[0]
    tp.text = "银发健康助手"
    style_paragraph(tp, size=52, bold=True, color=TITLE)
    tp.alignment = PP_ALIGN.CENTER

    sub_box = s1.shapes.add_textbox(Inches(0.8), Inches(3.35), Inches(8.4), Inches(0.8))
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = "面向65+中老年人的对话式健康管理应用"
    style_paragraph(sp, size=24, color=SUBTEXT)
    sp.alignment = PP_ALIGN.CENTER

    tag = s1.shapes.add_shape(1, Inches(2.6), Inches(4.35), Inches(4.8), Inches(0.58))
    tag.fill.solid()
    tag.fill.fore_color.rgb = HEADER_LINE
    tag.line.color.rgb = ACCENT
    tag.line.width = Pt(1)
    tag_tf = tag.text_frame
    tag_tf.text = "说话就能记 · 记得住 · 家人也放心"
    tag_p = tag_tf.paragraphs[0]
    style_paragraph(tag_p, size=20, bold=True, color=TITLE)
    tag_p.alignment = PP_ALIGN.CENTER
    tag_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    info_box = s1.shapes.add_textbox(Inches(0.8), Inches(5.35), Inches(8.4), Inches(1.3))
    info_tf = info_box.text_frame
    for i, line in enumerate([
        "智能应用项目开发综合实践 · 期末汇报",
        "202306120105  方晓楠",
        "2026年6月",
    ]):
        p = info_tf.paragraphs[0] if i == 0 else info_tf.add_paragraph()
        p.text = line
        style_paragraph(p, size=22, color=SUBTEXT)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

    # --- Slide 2 ---
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2, PAGE_BG)
    add_header_bar(s2, "① 目标用户与痛点  ·  约30秒")
    add_bullets(s2, [
        ("主要用户：65+独居/慢病老人", "高血压、糖尿病等需长期用药，智能手机操作能力有限"),
        ("辅助用户：异地子女（照护者）", "工作繁忙，无法实时掌握父母健康状况"),
        ("核心痛点", None),
        ("• 用药易忘漏服，纸笔记录难坚持", None),
        ("• 血压/血糖数据分散，看不懂趋势", None),
        ("• 突发不适时求助不便，子女不知情", None),
        ("• 通用App菜单复杂，老人不会找功能", None),
    ], font_size=22)

    # --- Slide 3 ---
    s3 = prs.slides.add_slide(blank)
    set_slide_bg(s3, PAGE_BG)
    add_header_bar(s3, "② 核心功能现场演示  ·  约1分30秒")

    tip = s3.shapes.add_textbox(Inches(0.45), Inches(1.05), Inches(9.1), Inches(0.55))
    tip_tf = tip.text_frame.paragraphs[0]
    tip_tf.text = "重点：说一句话 → 系统帮你做事（写入数据 / 触发提醒 / 跨会话记忆）"
    style_paragraph(tip_tf, size=20, bold=True, color=ACCENT_WARM, italic=True)

    add_demo_card(s3, 1.72, "演示1 · 聊天即操作", "我今天血压145/92",
                  "自动识别 → 写入血压档案 → 异常检测告警 → 子女端同步")
    add_demo_card(s3, 3.58, "演示2 · 长期记忆", "（第1天）我72岁，膝盖不好，有高血压",
                  "刷新页面后问「推荐运动」→ 结合年龄+膝盖+慢病个性化建议")
    add_demo_card(s3, 5.44, "演示3 · 主动服务 + 家庭闭环", "（设置）每天8点提醒吃降压药",
                  "到点推送提醒 → 确认/漏服记录 → 依从性统计 → 子女端查看")

    # --- Slide 4 ---
    s4 = prs.slides.add_slide(blank)
    set_slide_bg(s4, PAGE_BG)
    add_header_bar(s4, "③ 差异化亮点  ·  约30秒")

    table_shape = s4.shapes.add_table(5, 3, Inches(0.35), Inches(1.2), Inches(9.3), Inches(5.45))
    table = table_shape.table
    headers = ["对比维度", "普通方案", "银发健康助手"]
    rows = [
        ["交互方式", "找菜单填表单", "自然语言说话即可记录"],
        ["智能深度", "通用聊天问答", "三级意图识别 + 安全检查"],
        ["家庭关怀", "各自独立App", "用药追踪 → 子女端实时查看"],
        ["领域特色", "无慢病适配", "禁忌运动拦截 + 情绪关怀"],
    ]

    for ci, w in enumerate([Inches(2.1), Inches(3.4), Inches(3.8)]):
        table.columns[ci].width = w

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        for p in cell.text_frame.paragraphs:
            style_paragraph(p, size=22, bold=True, color=TITLE)
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                style_paragraph(p, size=20, bold=(ci == 0), color=TEXT if ci != 2 else ACCENT)
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ci == 2:
                cell.fill.solid()
                cell.fill.fore_color.rgb = HIGHLIGHT_BG
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT

    # --- Slide 5 ---
    s5 = prs.slides.add_slide(blank)
    set_slide_bg(s5, PAGE_BG)
    add_header_bar(s5, "④ 用户验证与迭代  ·  约30秒")
    add_bullets(s5, [
        ("真实用户调研（2026.6）", "3名真实目标用户：社区邻居68岁、父亲72岁、异地子女29岁"),
        ("可用性测试", "核心任务成功率：80% → 迭代后 100%（10/10）"),
        ("典型迭代案例", None),
        ("• 语音录入失败 → 增加「请这样说」模板 → 复测成功", None),
        ("• 预警误触发「头晕」→ 收紧紧急规则 → 信任度提升", None),
        ("• 子女端入口隐蔽 → 首页增加路径说明 → 无需口头指导", None),
        ("结论：3/3 用户愿意在协助下继续使用", None),
    ], font_size=22)

    # --- Slide 6: Thank you (light warm) ---
    s6 = prs.slides.add_slide(blank)
    set_slide_bg(s6, COVER_BG)
    add_warm_decor(s6)

    thanks = s6.shapes.add_textbox(Inches(1), Inches(2.4), Inches(8), Inches(1.2))
    tp6 = thanks.text_frame.paragraphs[0]
    tp6.text = "谢谢聆听"
    style_paragraph(tp6, size=48, bold=True, color=TITLE)
    tp6.alignment = PP_ALIGN.CENTER

    demo = s6.shapes.add_textbox(Inches(1), Inches(3.9), Inches(8), Inches(1.8))
    dtf = demo.text_frame
    for i, line in enumerate([
        "接下来进行现场演示",
        "父母端：localhost:5001  |  子女端：/child.html",
    ]):
        p = dtf.paragraphs[0] if i == 0 else dtf.add_paragraph()
        p.text = line
        style_paragraph(p, size=22, color=SUBTEXT)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(10)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
